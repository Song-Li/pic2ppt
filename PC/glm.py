import os
import base64
import json
import re
import cv2
import numpy as np
import glob
import torch
from openai import OpenAI
from simple_lama_inpainting import SimpleLama
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR

# ==============================================================================
# 【配置区域】
# ==============================================================================
# 请填入你的智谱 API Key
ZHIPU_API_KEY = "38a24dd463fb4002ab9113e2d5f6033c.Pu6ddxsZvONOJR6J" 

# 【坐标微调参数】如果发现位置有偏移，可以在这里微调
OFFSET_X = 0  # X轴偏移（像素），负数向左，正数向右
OFFSET_Y = 0  # Y轴偏移（像素），负数向上，正数向下

# 环境配置
os.environ["KMP_DUPLICATE_LIB_OK"] = "TRUE"

class ZhipuImageToPPT:
    def __init__(self, api_key):
        print("="*40)
        print("[-] 正在初始化...")
        
        # 1. 初始化 Client
        self.client = OpenAI(
            api_key=api_key,
            base_url="https://open.bigmodel.cn/api/paas/v4/"
        )
        print("[-] 智谱 API 连接已配置 (GLM-4.6V-Flash)")

        # 2. 初始化 LaMa
        print("[-] 正在加载 LaMa 修复模型...")
        self.lama = SimpleLama()
        if torch.cuda.is_available():
            print("[√] LaMa 模型已启用 GPU 加速")
        else:
            print("[!] LaMa 模型将运行在 CPU 上")
            
        print("[√] 初始化完成!")
        print("="*40)

    def encode_image_optimized(self, image_path):
        """优化的图片编码"""
        img = cv2.imdecode(np.fromfile(image_path, dtype=np.uint8), cv2.IMREAD_COLOR)
        if img is None: return None
        
        h, w = img.shape[:2]
        
        # 初始参数
        quality = 95
        max_size_bytes = 4 * 1024 * 1024  # 限制在 4MB 以内
        encoded_data = None
        
        # 循环压缩质量，直到体积达标
        while quality > 10:
            _, buffer = cv2.imencode(".jpg", img, [cv2.IMWRITE_JPEG_QUALITY, quality])
            if len(buffer) < max_size_bytes:
                encoded_data = buffer
                break
            quality -= 10

        # 如果质量降到 10 还是太大，缩小尺寸
        if encoded_data is None:
            print("    [Warn] 图片极大，被迫缩小尺寸...")
            scale = 0.5
            new_w, new_h = int(w * scale), int(h * scale)
            resized_img = cv2.resize(img, (new_w, new_h), interpolation=cv2.INTER_AREA)
            _, buffer = cv2.imencode(".jpg", resized_img, [cv2.IMWRITE_JPEG_QUALITY, 60])
            encoded_data = buffer

        return base64.b64encode(encoded_data).decode("utf-8").replace('\n', '')

    def get_text_from_zhipu(self, image_path, img_width, img_height, retry_count=0):
        """调用 GLM-4.6V-Flash 获取文字坐标和样式信息"""
        base64_image = self.encode_image_optimized(image_path)
        if not base64_image:
            print("[!] 图片读取失败")
            return []
        
        # ======================================================
        # 【优化的Prompt】使用0-1归一化坐标，更简洁准确
        # ======================================================
        if retry_count == 0:
            user_prompt = f"""
你是专业的OCR助手。请检测图片中所有文字的精确位置和样式。

图片尺寸：{img_width} x {img_height} 像素

【输出要求】仅输出JSON数组：
[
    {{
        "text": "文字内容",
        "box": [x_min, y_min, x_max, y_max],
        "size": 相对字号(0-100),
        "font": "字体名称",
        "bold": true/false,
        "italic": true/false
    }}
]

【坐标说明】
- box使用归一化坐标（0.0-1.0）
- 左上角是(0.0, 0.0)，右下角是(1.0, 1.0)
- 例如：图片中心的文字框可能是[0.4, 0.45, 0.6, 0.55]

【字号说明】
- size是相对值：大标题80-100，小标题50-70，正文25-45，注释10-20

【字体说明】
- 尽力识别字体：宋体、黑体、微软雅黑、Arial等
- 无法确定时返回null
"""
        else:
            # 简化重试
            user_prompt = f"""
检测图片中所有文字的位置。图片尺寸：{img_width}x{img_height}

输出JSON数组：
[{{"text": "文字", "box": [x_min, y_min, x_max, y_max]}}]

坐标使用0.0-1.0归一化值，左上角(0,0)，右下角(1,1)
"""

        try:
            response = self.client.chat.completions.create(
                model="glm-4.6v-flash",
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}},
                            {"type": "text", "text": user_prompt}
                        ]
                    }
                ],
                temperature=0.01,
                top_p=0.1,
                max_tokens=8000,
                stream=False
            )
            
            content = response.choices[0].message.content
            finish_reason = response.choices[0].finish_reason
            
            print("\n" + "="*80)
            print("【GLM-4.6V-Flash 原始返回结果】")
            print(f"模型: glm-4.6v-flash | 完成原因: {finish_reason}")
            if finish_reason == "length":
                print("⚠️  警告：输出因长度限制被截断！")
            print("="*80)
            print(content[:2000] + ("..." if len(content) > 2000 else ""))  # 只显示前2000字符
            print("="*80 + "\n")
            
            # 如果被截断，尝试修复JSON
            if finish_reason == "length":
                print("    [!] 检测到截断，尝试修复JSON...")
                content = self._fix_truncated_json(content)
            
            clean_content = re.sub(r'^```json\s*|\s*```$', '', content.strip(), flags=re.IGNORECASE)

            try:
                data = json.loads(clean_content)
                if isinstance(data, list):
                    print(f"    [√] JSON解析成功，共 {len(data)} 个文字块")
                    return data
                elif isinstance(data, dict):
                    items = data.get('items', []) or data.get('words', []) or data.get('texts', [])
                    print(f"    [√] JSON解析成功，共 {len(items)} 个文字块")
                    return items
                return data
            except json.JSONDecodeError as e:
                print(f"    [Warn] JSON解析失败: {e}")
                
                if retry_count == 0:
                    print("    [>] 使用简化Prompt重试...")
                    return self.get_text_from_zhipu(image_path, img_width, img_height, retry_count=1)
                
                print("    [>] 尝试正则提取...")
                items = self._regex_extract(content)
                if items:
                    print(f"    [√] 正则提取成功: {len(items)} 个文字块")
                return items

        except Exception as e:
            import traceback
            print(f"    [API 异常] {e}")
            traceback.print_exc()
            return []

    def _fix_truncated_json(self, content):
        """修复被截断的JSON"""
        last_complete = content.rfind('}')
        if last_complete != -1:
            content = content[:last_complete+1]
            if not content.rstrip().endswith(']'):
                content = content.rstrip().rstrip(',') + '\n]'
        return content

    def _regex_extract(self, content):
        """正则表达式提取文字块"""
        items = []
        
        # 匹配归一化坐标（0.0-1.0）
        pattern1 = r'\{\s*"text":\s*"([^"]*)",\s*"box":\s*\[([\d.]+),\s*([\d.]+),\s*([\d.]+),\s*([\d.]+)\](?:,\s*"size":\s*([\d.]+))?(?:,\s*"font":\s*"([^"]*)")?(?:,\s*"bold":\s*(true|false))?(?:,\s*"italic":\s*(true|false))?\s*\}'
        
        matches = re.finditer(pattern1, content)
        for m in matches:
            item = {
                "text": m.group(1),
                "box": [float(m.group(2)), float(m.group(3)), float(m.group(4)), float(m.group(5))]
            }
            if m.group(6):
                item["size"] = float(m.group(6))
            if m.group(7):
                item["font"] = m.group(7)
            if m.group(8):
                item["bold"] = m.group(8) == "true"
            if m.group(9):
                item["italic"] = m.group(9) == "true"
            items.append(item)
        
        return items

    def process_folder(self, folder_path, output_pptx_path):
        debug_folder = "debug_output_zhipu"
        if not os.path.exists(debug_folder):
            os.makedirs(debug_folder)

        image_files = []
        for ext in ['*.jpg', '*.jpeg', '*.png', '*.bmp']:
            image_files.extend(glob.glob(os.path.join(folder_path, ext)))
        image_files.sort()
        
        if not image_files:
            print(f"[!] 错误: 文件夹 '{folder_path}' 中未找到图片。")
            return

        print(f"[-] 找到 {len(image_files)} 张图片，准备处理...")
        print(f"[-] 坐标微调参数: X偏移={OFFSET_X}px, Y偏移={OFFSET_Y}px")
        
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9
        prs.slide_height = Inches(7.5)

        for idx, img_path in enumerate(image_files):
            filename = os.path.basename(img_path)
            print(f"\n>>> [{idx+1}/{len(image_files)}] 正在处理: {filename}")
            try:
                self.process_single_slide(prs, img_path, debug_folder)
            except Exception as e:
                import traceback
                print(f"[!] 处理失败: {e}")
                traceback.print_exc()
                continue

        prs.save(output_pptx_path)
        print(f"\n[√] 全部完成! PPT 已保存至: {output_pptx_path}")

    def process_single_slide(self, prs, image_path, debug_folder):
        # 1. 读取原始图片
        img_original = cv2.imdecode(np.fromfile(image_path, dtype=np.uint8), cv2.IMREAD_COLOR)
        if img_original is None:
            print(f"    [!] 无法读取图片: {image_path}")
            return
            
        h, w = img_original.shape[:2]
        print(f"    [图片尺寸] {w} x {h}")

        # 2. 调用 API
        api_data = self.get_text_from_zhipu(image_path, w, h)
        
        text_data = []
        
        if api_data:
            print(f"    [成功] 获取 {len(api_data)} 个文字块")
            
            for i, item in enumerate(api_data):
                content = item.get('text', '').strip()
                box_norm = item.get('box', [])
                
                if not content or len(box_norm) != 4:
                    print(f"      [跳过{i+1}] 数据不完整")
                    continue
                
                try:
                    # 【关键】归一化坐标直接转换为像素坐标
                    x_min_norm, y_min_norm, x_max_norm, y_max_norm = [float(x) for x in box_norm]
                    
                    # 检查坐标范围
                    if max(x_min_norm, y_min_norm, x_max_norm, y_max_norm) > 10:
                        # 可能是1000坐标系，转换为0-1
                        x_min_norm /= 1000
                        y_min_norm /= 1000
                        x_max_norm /= 1000
                        y_max_norm /= 1000
                        print(f"      [注意] 检测到1000坐标系，已自动转换")
                    
                    # 裁剪到有效范围
                    x_min_norm = np.clip(x_min_norm, 0, 1)
                    y_min_norm = np.clip(y_min_norm, 0, 1)
                    x_max_norm = np.clip(x_max_norm, 0, 1)
                    y_max_norm = np.clip(y_max_norm, 0, 1)
                    
                    # 转换为像素坐标
                    x1 = int(x_min_norm * w) + OFFSET_X
                    y1 = int(y_min_norm * h) + OFFSET_Y
                    x2 = int(x_max_norm * w) + OFFSET_X
                    y2 = int(y_max_norm * h) + OFFSET_Y
                    
                    # 确保有效
                    x1, x2 = max(0, min(x1, x2)), max(x1, x2, x1+5)
                    y1, y2 = max(0, min(y1, y2)), max(y1, y2, y1+5)
                    
                    rect_w = x2 - x1
                    rect_h = y2 - y1
                    
                    print(f"      [{i+1}] \"{content[:20]}\" | 归一化:[{x_min_norm:.3f},{y_min_norm:.3f},{x_max_norm:.3f},{y_max_norm:.3f}] | 像素:[{x1},{y1},{x2},{y2}]")
                    
                except Exception as e:
                    print(f"      [跳过{i+1}] 坐标解析失败: {e}")
                    continue
                
                # 颜色提取
                pad = 3
                safe_x = np.clip(x1 + pad, 0, w - 1)
                safe_y = np.clip(y1 + pad, 0, h - 1)
                safe_w = np.clip(rect_w - 2*pad, 1, w - safe_x)
                safe_h = np.clip(rect_h - 2*pad, 1, h - safe_y)
                
                roi = img_original[safe_y:safe_y+safe_h, safe_x:safe_x+safe_w]
                
                final_color = (0, 0, 0)
                if roi.size > 0:
                    try:
                        gray_roi = cv2.cvtColor(roi, cv2.COLOR_BGR2GRAY)
                        min_val, _, min_loc, _ = cv2.minMaxLoc(gray_roi)
                        b, g, r = roi[min_loc[1], min_loc[0]]
                        brightness = 0.299*r + 0.587*g + 0.114*b
                        if brightness > 220:
                            final_color = (0, 0, 0)
                        else:
                            final_color = (int(r), int(g), int(b))
                    except:
                        final_color = (0, 0, 0)

                # 字体信息
                font_size_norm = item.get('size', None)
                font_family = item.get('font', None)
                bold = item.get('bold', False)
                italic = item.get('italic', False)
                
                # 【改进字号计算】
                if font_size_norm and font_size_norm > 0:
                    # size是0-100，映射到合理的像素范围
                    # 假设：100对应80pt，即约107px（在1920宽的图上）
                    # 缩放到实际图片尺寸
                    base_size = (font_size_norm / 100) * (h / 1080) * 80  # 基于1080p标准
                    actual_font_size_px = base_size
                else:
                    # 根据框高估算
                    actual_font_size_px = rect_h * 0.85

                text_data.append({
                    "text": content,
                    "rect": (x1, y1, rect_w, rect_h),
                    "color": final_color,
                    "font_size_px": actual_font_size_px,
                    "font_family": font_family,
                    "bold": bold,
                    "italic": italic
                })
                
                display_text = content[:15] + "..." if len(content) > 15 else content
                print(f"        文字:\"{display_text}\" | 字号:{actual_font_size_px:.1f}px | 字体:{font_family} | 粗:{bold}")

        else:
            print("    [警告] API 返回空数据")
            return

        # 3. 绘制 Debug 图
        print(f"    [Debug] 绘制调试图...")
        debug_img = img_original.copy()
        
        for idx, item in enumerate(text_data):
            x, y, rw, rh = item['rect']
            
            # 绘制红色边框（加粗）
            cv2.rectangle(debug_img, (x, y), (x+rw, y+rh), (0, 0, 255), 4)
            
            # 绘制序号标签（带背景）
            label = str(idx+1)
            font_scale = 1.0
            thickness = 2
            (label_w, label_h), baseline = cv2.getTextSize(label, cv2.FONT_HERSHEY_SIMPLEX, font_scale, thickness)
            
            # 绿色背景
            cv2.rectangle(debug_img, 
                         (x, y - label_h - 10), 
                         (x + label_w + 10, y), 
                         (0, 255, 0), -1)
            
            # 黑色文字
            cv2.putText(debug_img, label, 
                       (x + 5, y - 5), 
                       cv2.FONT_HERSHEY_SIMPLEX, 
                       font_scale, (0, 0, 0), thickness)
        
        debug_path = os.path.join(debug_folder, f"debug_{os.path.basename(image_path)}")
        cv2.imencode('.png', debug_img)[1].tofile(debug_path)
        print(f"    [√] Debug图已保存: {debug_path}")

        # 4. LaMa 抹除
        if text_data:
            print("    [LaMa] 移除背景文字...")
            mask = np.zeros((h, w), dtype=np.uint8)
            for item in text_data:
                x, y, rw, rh = item['rect']
                pad_x, pad_y = int(rw * 0.1), int(rh * 0.1)
                cv2.rectangle(mask, 
                            (max(0, x-pad_x), max(0, y-pad_y)), 
                            (min(w, x+rw+pad_x), min(h, y+rh+pad_y)), 
                            255, -1)
            
            kernel = np.ones((9, 9), np.uint8)
            mask = cv2.dilate(mask, kernel, iterations=2)
            
            img_rgb = cv2.cvtColor(img_original, cv2.COLOR_BGR2RGB)
            clean_img_pil = self.lama(img_rgb, mask)
            clean_img_rgb = np.array(clean_img_pil)
            clean_img_bgr = cv2.cvtColor(clean_img_rgb, cv2.COLOR_RGB2BGR)
        else:
            clean_img_bgr = img_original

        # 5. 生成 PPT
        ppt_width_px = int(prs.slide_width / 914400 * 96 * 10)
        ppt_height_px = int(prs.slide_height / 914400 * 96 * 10)
        
        clean_img_resized = cv2.resize(clean_img_bgr, (ppt_width_px, ppt_height_px), 
                                       interpolation=cv2.INTER_LANCZOS4)
        
        temp_bg_filename = f"temp_bg_{os.path.basename(image_path)}.png"
        cv2.imencode('.png', clean_img_resized)[1].tofile(temp_bg_filename)

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(temp_bg_filename, 0, 0, 
                                width=prs.slide_width, height=prs.slide_height)
        
        # 6. 添加文本框
        scale_x = prs.slide_width / w
        scale_y = prs.slide_height / h
        
        print(f"    [PPT] 缩放比例 X:{scale_x/914400:.6f}, Y:{scale_y/914400:.6f}")
        
        for idx, item in enumerate(text_data):
            x, y, rw, rh = item['rect']
            
            # 转换为EMU
            left_emu = int(x * scale_x)
            top_emu = int(y * scale_y)
            width_emu = int(rw * scale_x)
            height_emu = int(rh * scale_y)
            
            # 添加文本框
            shape = slide.shapes.add_textbox(left_emu, top_emu, width_emu, height_emu)
            tf = shape.text_frame
            tf.word_wrap = True
            tf.margin_left = 0
            tf.margin_right = 0
            tf.margin_top = 0
            tf.margin_bottom = 0
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
            
            p = tf.paragraphs[0]
            p.text = item['text']
            
            # 颜色
            c = item['color']
            p.font.color.rgb = RGBColor(c[0], c[1], c[2])
            
            # 【改进字号计算】
            font_size_px = item['font_size_px']
            
            # 字号转换：像素 -> 点数
            # PPT中1英寸 = 72点 = 96像素（屏幕DPI）
            # 但要考虑图片到PPT的缩放
            img_dpi = 96
            ppt_inch_per_px = 1 / 914400  # EMU转英寸
            
            # 图片像素对应的PPT英寸
            font_inch = (font_size_px / img_dpi) * (prs.slide_height / 914400) / (h / img_dpi)
            font_pt = font_inch * 72
            
            # 或者更简单的方法：直接按比例缩放
            font_pt = font_size_px * (prs.slide_height / 914400 / (h / 96)) * 72 / 96
            
            # 限制范围
            font_pt = max(8, min(72, font_pt))
            
            try:
                p.font.size = Pt(int(font_pt))
                print(f"      [{idx+1}] 字号: {font_size_px:.1f}px -> {font_pt:.1f}pt")
            except Exception as e:
                print(f"      [{idx+1}] 字号设置失败，使用默认12pt")
                p.font.size = Pt(12)
            
            # 字体
            if item.get('font_family'):
                font_map = {
                    '宋体': 'SimSun', '黑体': 'SimHei', '微软雅黑': 'Microsoft YaHei',
                    '楷体': 'KaiTi', '仿宋': 'FangSong', 'Arial': 'Arial',
                    'Times': 'Times New Roman', 'Calibri': 'Calibri'
                }
                
                font_name = item['font_family']
                for key, value in font_map.items():
                    if key.lower() in font_name.lower():
                        font_name = value
                        break
                
                try:
                    p.font.name = font_name
                except:
                    pass
            
            # 粗体斜体
            if item.get('bold'):
                p.font.bold = True
            if item.get('italic'):
                p.font.italic = True

        if os.path.exists(temp_bg_filename):
            os.remove(temp_bg_filename)
        
        print(f"    [√] 幻灯片完成，{len(text_data)} 个文本框")

if __name__ == "__main__":
    if "你的API_KEY" in ZHIPU_API_KEY:
        print("[!] 请先填入 API Key！")
    else:
        input_folder = "images"
        output_file = "result_zhipu_fix.pptx"
        
        if not os.path.exists(input_folder):
            os.makedirs(input_folder)
            print(f"[!] 请在 '{input_folder}' 文件夹中放入图片")
        else:
            print("\n【提示】如果生成的PPT文字位置有偏移，可以调整代码顶部的 OFFSET_X 和 OFFSET_Y 参数")
            print("例如：OFFSET_X = -2  # 向左偏移2像素")
            print("     OFFSET_Y = 3   # 向下偏移3像素\n")
            
            converter = ZhipuImageToPPT(ZHIPU_API_KEY)
            converter.process_folder(input_folder, output_file)
