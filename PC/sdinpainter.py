import os
import cv2
import numpy as np
import glob
import torch
from PIL import Image
from paddleocr import PaddleOCR
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from diffusers import StableDiffusionInpaintPipeline

# ================= 配置区域 =================
INPUT_FOLDER = "images"           # 图片文件夹
OUTPUT_PPT = "sd_output.pptx"     # 输出PPT名称
MIN_SCORE = 0.60                  # OCR置信度阈值
FILTER_ICON_RATIO = 1.2           # 图标过滤阈值
USE_FP16 = True                   # 是否开启半精度加速 (显存<8G请改为False)
# ===========================================

class SDInpainter:
    def __init__(self):
        print("[-] 正在加载 Stable Diffusion Inpainting 模型 (需下载约4GB，请耐心等待)...")
        self.model_id = "runwayml/stable-diffusion-inpainting"
        
        # 加载配置
        dtype = torch.float16 if USE_FP16 else torch.float32
        try:
            self.pipe = StableDiffusionInpaintPipeline.from_pretrained(
                self.model_id,
                torch_dtype=dtype,
                variant="fp16" if USE_FP16 else None
            )
            
            if torch.cuda.is_available():
                self.pipe.to("cuda")
                print("[√] SD Inpainting 已加载至 GPU")
            elif torch.backends.mps.is_available(): # Mac M1/M2
                self.pipe.to("mps")
                print("[√] SD Inpainting 已加载至 MPS (Mac)")
            else:
                print("[!] 警告: 未检测到 GPU，CPU 运行 SD 可能会非常慢！")
                self.pipe.to("cpu")
            
            # 开启显存优化
            self.pipe.enable_attention_slicing()
            
        except Exception as e:
            print(f"[!] 模型加载失败: {e}")
            raise e

    def __call__(self, image_pil, mask_pil):
        """
        使得类实例可以直接被调用，解决 'object is not callable' 问题
        image_pil: RGB 格式的 PIL Image
        mask_pil:  单通道 PIL Image (白色为抹除区域)
        """
        # SD 接受的图片宽高最好是 8 的倍数，这里不做强制缩放，依赖 diffusers 内部处理
        # Prompt 设为空，让模型根据背景自动推断填充内容
        with torch.no_grad():
            output = self.pipe(
                prompt="background, clean, high quality", 
                negative_prompt="text, watermark, writing, artifacts, messy, blur, low quality",
                image=image_pil,
                mask_image=mask_pil,
                num_inference_steps=25, # 步数：20-30 之间速度与质量平衡最好
                guidance_scale=7.5
            ).images[0]
        return output

class ImageToPPTConverter:
    def __init__(self):
        # 1. 初始化 SD
        self.sd_model = SDInpainter()
        
        # 2. 初始化 OCR
        print("[-] 正在初始化 PaddleOCR...")
        self.ocr = PaddleOCR(use_angle_cls=True, lang="ch", show_log=False)
        print("[√] 初始化完成")

    def get_dominant_color(self, img, box):
        """提取文字颜色"""
        try:
            x, y, w, h = box
            # 取中心的一小块区域
            crop = img[y+h//4 : y+3*h//4, x+w//4 : x+3*w//4]
            if crop.size == 0: return (0, 0, 0)
            
            # 计算均值颜色
            b, g, r = cv2.mean(crop)[:3]
            
            # 简单判断：如果背景非常亮，颜色可能是黑色；反之亦然
            # 这里直接返回检测到的颜色，如果觉得不准可以强制设为黑色 (0,0,0)
            return (int(r), int(g), int(b))
        except:
            return (0, 0, 0)

    def is_likely_icon(self, w, h, text):
        """过滤掉可能是图标的小方块"""
        ratio = max(w, h) / min(w, h)
        # 尺寸小 + 形状方正 + 字数少 = 极有可能是图标
        if w < 60 and h < 60 and ratio < FILTER_ICON_RATIO and len(text) < 2:
            return True
        return False

    def process_folder(self, folder_path, output_path):
        image_files = sorted(glob.glob(os.path.join(folder_path, "*.[jpg][png]*")))
        
        if not image_files:
            print(f"[!] 文件夹 {folder_path} 为空")
            return

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for i, img_path in enumerate(image_files):
            print(f"\n>>> [{i+1}/{len(image_files)}] 正在处理: {os.path.basename(img_path)}")
            self.process_one_slide(prs, img_path)

        prs.save(output_path)
        print(f"\n[√] 处理完成，PPT已保存: {output_path}")

    def process_one_slide(self, prs, img_path):
        # 读取图片
        img_bgr = cv2.imdecode(np.fromfile(img_path, dtype=np.uint8), cv2.IMREAD_COLOR)
        h, w = img_bgr.shape[:2]

        # 1. OCR 识别
        result = self.ocr.ocr(img_bgr, cls=True)
        
        valid_blocks = []
        # 创建 Mask，黑色背景，白色为需要抹除的区域
        mask = np.zeros((h, w), dtype=np.uint8)

        if result and result[0]:
            for line in result[0]:
                coords = line[0]
                text = line[1][0]
                score = line[1][1]

                if score < MIN_SCORE: continue
                
                coords_np = np.array(coords, dtype=np.int32)
                x, y, rw, rh = cv2.boundingRect(coords_np)
                
                # 过滤图标
                if self.is_likely_icon(rw, rh, text):
                    print(f"    [忽略图标] {text}")
                    continue

                color = self.get_dominant_color(img_bgr, (x, y, rw, rh))
                
                valid_blocks.append({
                    "text": text,
                    "rect": (x, y, rw, rh),
                    "color": color
                })

                # 在 Mask 上画白色实心矩形
                cv2.fillPoly(mask, [coords_np], 255)

        # 2. SD 抹除文字
        if valid_blocks:
            print(f"    [SD] 正在重绘背景 (共 {len(valid_blocks)} 处文字)...")
            
            # 膨胀 Mask：SD 需要更多的上下文边缘，稍微多膨胀一点效果更好
            kernel = np.ones((8, 8), np.uint8)
            mask_dilated = cv2.dilate(mask, kernel, iterations=2) # 膨胀2次
            
            # 转换格式给 SD
            img_rgb = cv2.cvtColor(img_bgr, cv2.COLOR_BGR2RGB)
            img_pil = Image.fromarray(img_rgb)
            mask_pil = Image.fromarray(mask_dilated)
            
            # === 关键修正：直接调用类实例 ===
            clean_pil = self.sd_model(img_pil, mask_pil)
            
            clean_bgr = cv2.cvtColor(np.array(clean_pil), cv2.COLOR_RGB2BGR)
        else:
            clean_bgr = img_bgr

        # 3. 生成 PPT
        temp_bg = f"temp_bg_{os.path.basename(img_path)}.png"
        cv2.imencode('.png', clean_bgr)[1].tofile(temp_bg)
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(temp_bg, 0, 0, width=prs.slide_width, height=prs.slide_height)
        
        # 4. 回填文字
        scale_x = prs.slide_width / w
        scale_y = prs.slide_height / h
        
        for block in valid_blocks:
            x, y, rw, rh = block['rect']
            
            left = int(x * scale_x)
            top = int(y * scale_y)
            width = int(rw * scale_x)
            height = int(rh * scale_y)
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            tf = textbox.text_frame
            tf.word_wrap = False
            tf.margin_top = 0
            tf.margin_left = 0
            
            p = tf.paragraphs[0]
            p.text = block['text']
            # p.font.color.rgb = RGBColor(*block['color']) # 使用提取颜色
            p.font.color.rgb = RGBColor(0, 0, 0) # 还是建议默认黑色/深色，提取的颜色有时候会太浅看不清
            
            # 动态字号
            # 72磅 = 1英寸。假设文字高度占框高的 70%
            # rh (px) -> inches -> pt
            font_size_pt = (rh / h) * prs.slide_height.inches / 914400 * 72 * 0.7 
            # 修正计算：inches本身是emu单位，直接转
            font_size_pt = (rh * scale_y) / 12700 * 0.7 # 粗略换算
            
            p.font.size = Pt(max(10, min(60, font_size_pt)))
            p.font.name = "微软雅黑"

        if os.path.exists(temp_bg):
            os.remove(temp_bg)

if __name__ == "__main__":
    if not os.path.exists(INPUT_FOLDER):
        os.makedirs(INPUT_FOLDER)
        print(f"[!] 请先创建 {INPUT_FOLDER} 文件夹并放入图片")
    else:
        try:
            converter = ImageToPPTConverter()
            converter.process_folder(INPUT_FOLDER, OUTPUT_PPT)
        except Exception as e:
            print(f"[!] 程序崩溃: {e}")
            import traceback
            traceback.print_exc()