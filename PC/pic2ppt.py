import os
import sys
import cv2
import numpy as np
import glob
from PIL import Image
import torch
from simple_lama_inpainting import SimpleLama
from paddleocr import PaddleOCR
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_VERTICAL_ANCHOR

# ================= 顶级配置 (最终版 - 颜色/Mask修正) =================
INPUT_FOLDER = "images"
OUTPUT_PPT = "result_ultimate_final.pptx"
DEBUG_FOLDER = "debug_ultimate_final"

# 1. 模型路径 (Server模型配置)
MODEL_DIR = "models_server"
DET_DIR = os.path.join(MODEL_DIR, "ch_PP-OCRv4_det_server_infer")
REC_DIR = os.path.join(MODEL_DIR, "ch_PP-OCRv4_rec_server_infer")
CLS_DIR = os.path.join(MODEL_DIR, "ch_ppocr_mobile_v2.0_cls_infer")

# 2. 调优参数
MIN_SCORE = 0.75             
FONT_SIZE_FACTOR = 0.98      
MASK_PADDING_PIXELS = 0      # [修正] Mask安全边距降至1像素
MASK_DILATE_ITER = 1         # [修正] 核心：Mask 膨胀迭代次数降至0 (仅使用 kernel size)
PADDING_PIXELS_GREEN = 0     # 绿框的固定外扩像素 
# ===========================================

def check_gpu():
    print("[-] 正在检查 GPU 环境...")
    if not torch.cuda.is_available():
        print("[!] 警告：PyTorch 未检测到 GPU！程序将以 CPU 运行。")
    else:
        name = torch.cuda.get_device_name(0)
        print(f"[√] 检测到显卡: {name} (CUDA可用)")

class ImageToPPTConverter:
    def __init__(self):
        if not os.path.exists(DET_DIR):
            raise FileNotFoundError(f"找不到模型文件: {DET_DIR}\n请先运行 download_models.py")
            
        print("[-] 加载 PaddleOCR Server 模型 (GPU加速中)...")
        self.ocr = PaddleOCR(
            text_detection_model_dir=DET_DIR,
            text_recognition_model_dir=REC_DIR,
            textline_orientation_model_dir=CLS_DIR,
            use_textline_orientation=True,
            device="gpu" 
        )
        
        print("[-] 加载 LaMa 修复模型 (GPU加速中)...")
        self.lama = SimpleLama()
        print("[√] 初始化完成")

    def is_noise_or_icon(self, w, h, text, score):
        if w < 10 or h < 10: return True
        ratio = w / h
        if 0.8 < ratio < 1.3 and max(w, h) < 60 and len(text) < 2:
            if '\u4e00' <= text <= '\u9fff': return False 
            return True
        if len(text) == 1 and not ('\u4e00' <= text <= '\u9fff' or text.isalnum()):
            return True
        return False

    def smart_refine_box(self, img, x, y, w, h):
        if w <= 0 or h <= 0: return x, y, w, h
        
        pad_probe = 2 
        H, W = img.shape[:2]
        y1, y2 = max(0, y - pad_probe), min(H, y + h + pad_probe)
        x1, x2 = max(0, x - pad_probe), min(W, x + w + pad_probe)
        roi = img[y1:y2, x1:x2]
        
        if roi.size == 0: return x, y, w, h

        gray = cv2.cvtColor(roi, cv2.COLOR_BGR2GRAY)
        blurred = cv2.GaussianBlur(gray, (3, 3), 0)
        _, binary = cv2.threshold(blurred, 0, 255, cv2.THRESH_BINARY_INV + cv2.THRESH_OTSU)
        
        points = cv2.findNonZero(binary)
        
        if points is not None:
            rx, ry, rw, rh = cv2.boundingRect(points)
            
            real_x = x1 + rx
            real_y = y1 + ry
            
            final_x = max(0, real_x - PADDING_PIXELS_GREEN)
            final_y = max(0, real_y - PADDING_PIXELS_GREEN)
            final_w = min(W - final_x, rw + 2*PADDING_PIXELS_GREEN)
            final_h = min(H - final_y, rh + 2*PADDING_PIXELS_GREEN)
            
            return int(final_x), int(final_y), int(final_w), int(final_h)
            
        return x, y, w, h

    def analyze_style(self, img, box):
        """
        [核心修复] 颜色和粗体分析：基于 HSV 空间，提取最饱和的颜色作为字体色
        """
        x, y, w, h = box
        if w <= 0 or h <= 0: return (0,0,0), False
        
        roi = img[y:y+h, x:x+w] 
        if roi.size == 0: return (0,0,0), False

        is_bold = False
        try:
            # 1. 转换为 HSV 空间
            roi_hsv = cv2.cvtColor(roi, cv2.COLOR_BGR2HSV)
            
            # 2. 创建一个阈值 Mask，过滤掉极暗/极亮的背景 (V<50 或 V>200)
            # 目标是排除纯黑/纯白背景，只保留彩色或中性色文字
            lower_bound = np.array([0, 0, 50]) # V > 50 (排除纯黑)
            upper_bound = np.array([179, 255, 200]) # V < 200 (排除纯白高光)
            hsv_mask = cv2.inRange(roi_hsv, lower_bound, upper_bound)
            
            # 3. 提取有效的彩色像素点
            valid_pixels = roi[hsv_mask != 0]
            
            if valid_pixels.size > 0:
                # 再次转换为 HSV，找到饱和度最高的点 (S值最大)
                valid_pixels_hsv = cv2.cvtColor(valid_pixels.reshape(1, -1, 3).astype(np.uint8), cv2.COLOR_BGR2HSV).reshape(-1, 3)
                
                # 找到 S (饱和度) 最大的像素的索引
                max_saturation_index = np.argmax(valid_pixels_hsv[:, 1])
                
                # 返回该点的原始 BGR 颜色
                color_bgr = valid_pixels[max_saturation_index]
            else:
                # 如果没有饱和度高的像素，回退到原始 BGR 空间的平均颜色
                color_bgr = cv2.mean(roi)[:3]

            text_color = (int(color_bgr[2]), int(color_bgr[1]), int(color_bgr[0]))
            
            # --- 粗体判断 ---
            gray = cv2.cvtColor(roi, cv2.COLOR_BGR2GRAY)
            _, binary = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY_INV + cv2.THRESH_OTSU)
            if (cv2.countNonZero(binary) / binary.size) > 0.45:
                is_bold = True
        except: 
            # 失败则回退到中心点采样
            try:
                sample = roi[max(0, h//2), max(0, w//2)] 
                text_color = (int(sample[2]), int(sample[1]), int(sample[0]))
            except:
                text_color = (0,0,0)

        return text_color, is_bold

    def process_folder(self, folder_path, output_path):
        if not os.path.exists(DEBUG_FOLDER): os.makedirs(DEBUG_FOLDER) 
        
        image_files = sorted(glob.glob(os.path.join(folder_path, "*.[jpg][png][jpeg]*")))
        
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for i, img_path in enumerate(image_files):
            print(f">>> [{i+1}/{len(image_files)}] 处理: {os.path.basename(img_path)}")
            try:
                self.process_one_slide(prs, img_path)
            except Exception as e:
                print(f"   [!] 错误: {e}")
                import traceback
                traceback.print_exc()

        prs.save(output_path)
        print(f"\n[√] 完成！PPT: {output_path}")

    def process_one_slide(self, prs, img_path):
        filename = os.path.splitext(os.path.basename(img_path))[0]
        img_bgr = cv2.imdecode(np.fromfile(img_path, dtype=np.uint8), cv2.IMREAD_COLOR)
        if img_bgr is None: return
        h_img, w_img = img_bgr.shape[:2]
        
        debug_img = img_bgr.copy()
        
        # 1. OCR (Server Model)
        result = self.ocr.ocr(img_bgr)
        
        valid_blocks = []
        mask = np.zeros((h_img, w_img), dtype=np.uint8)

        lines = result[0] if (result and result[0]) else []

        for line in lines:
            if not line: continue
            coords = line[0]
            text_content = line[1][0]
            score = line[1][1]
            
            if score < MIN_SCORE: continue

            coords_np = np.array(coords, dtype=np.int32)
            raw_x, raw_y, raw_w, raw_h = cv2.boundingRect(coords_np)
            
            # === [步骤 1] 强力过滤图标/噪点 (基于红框) ===
            if self.is_noise_or_icon(raw_w, raw_h, text_content, score):
                cv2.rectangle(debug_img, (raw_x, raw_y), (raw_x+raw_w, raw_y+raw_h), (128,128,128), 1)
                continue
            
            # === [步骤 2] 计算绿框 (精修尺寸) ===
            g_x, g_y, g_rw, g_rh = self.smart_refine_box(img_bgr, raw_x, raw_y, raw_w, raw_h)
            
            if g_rw < 5 or g_rh < 5: continue
            
            # === 3. 确定最终使用的框和数据 ===
            # PPT Textbox 定位、尺寸 (x, y, w, h) 全部使用绿框的尺寸
            x, y, rw, rh = g_x, g_y, g_rw, g_rh
            
            # 颜色提取使用绿框坐标
            color, is_bold = self.analyze_style(img_bgr, (x, y, rw, rh)) 
            
            # 绘制红框 (原始OCR框) 和绿框 (最终使用的框)
            cv2.rectangle(debug_img, (raw_x, raw_y), (raw_x+raw_w, raw_y+raw_h), (0,0,255), 1)
            cv2.rectangle(debug_img, (x, y), (x+rw, y+rh), (0,255,0), 2)

            valid_blocks.append({
                "text": text_content,
                "rect": (x, y, rw, rh), 
                "color": color,
                "bold": is_bold
            })
            
            # Mask 扩张：使用红框 (raw_x/raw_w) + 额外安全边距来擦除
            m_x1 = max(0, raw_x - MASK_PADDING_PIXELS)
            m_y1 = max(0, raw_y - MASK_PADDING_PIXELS)
            m_x2 = min(w_img, raw_x + raw_w + MASK_PADDING_PIXELS)
            m_y2 = min(h_img, raw_y + raw_h + MASK_PADDING_PIXELS)
            cv2.rectangle(mask, (m_x1, m_y1), (m_x2, m_y2), 255, -1)

        # 4. LaMa 修复
        if valid_blocks:
            cv2.imencode('.jpg', debug_img)[1].tofile(f"{DEBUG_FOLDER}/{filename}_final_boxes.jpg")
            
            kernel = np.ones((3, 3), np.uint8)
            # 核心：降低膨胀迭代次数
            mask_dilated = cv2.dilate(mask, kernel, iterations=MASK_DILATE_ITER)
            
            img_rgb = cv2.cvtColor(img_bgr, cv2.COLOR_BGR2RGB)
            try:
                clean_pil = self.lama(Image.fromarray(img_rgb), Image.fromarray(mask_dilated))
                clean_bgr = cv2.cvtColor(np.array(clean_pil), cv2.COLOR_RGB2BGR)
                cv2.imencode('.jpg', mask_dilated)[1].tofile(f"{DEBUG_FOLDER}/{filename}_mask.jpg")
                cv2.imencode('.jpg', clean_bgr)[1].tofile(f"{DEBUG_FOLDER}/{filename}_clean.jpg")
            except Exception as e:
                print(f"   [!] LaMa Error: {e}")
                clean_bgr = img_bgr
        else:
            clean_bgr = img_bgr

        # 5. 插入 PPT
        temp_bg = f"temp_{filename}.png"
        cv2.imencode('.png', clean_bgr)[1].tofile(temp_bg)
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(temp_bg, 0, 0, width=prs.slide_width, height=prs.slide_height)
        
        scale_x = prs.slide_width / w_img
        scale_y = prs.slide_height / h_img
        
        for block in valid_blocks:
            x, y, rw, rh = block['rect']
            text_content = block['text']
            
            textbox = slide.shapes.add_textbox(
                int(x * scale_x), int(y * scale_y), 
                int(rw * scale_x), int(rh * scale_y)
            )
            tf = textbox.text_frame
            tf.word_wrap = False 
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0
            
            p = tf.paragraphs[0]
            p.text = text_content
            p.font.color.rgb = RGBColor(*block['color'])
            p.font.bold = block['bold']
            p.font.name = "微软雅黑" # [要求] 字体统一使用微软雅黑
            
            # --- 字号双重约束 ---
            # 1. 高度约束 (H_pt): 基于绿框高度
            H_pt = (rh / h_img) * (prs.slide_height.inches * 72)
            
            # 2. 宽度约束 (W_pt_constraint): 确保字不超框
            char_count = len(text_content)
            W_pt = (rw / w_img) * (prs.slide_width.inches * 72)
            
            # W_pt_constraint = Textbox宽度 / 字符数 / 1.0 (保守系数，保证不超框)
            W_pt_constraint = W_pt / char_count / 1.0 
            
            # 3. 取两者中较小的值 (同时满足高宽限制)
            final_pt_size = min(H_pt, W_pt_constraint)
            
            p.font.size = Pt(max(9, final_pt_size * FONT_SIZE_FACTOR))
            # --------------------------------

        if os.path.exists(temp_bg): os.remove(temp_bg)

if __name__ == "__main__":
    check_gpu()
    if not os.path.exists(INPUT_FOLDER):
        os.makedirs(INPUT_FOLDER)
    else:
        converter = ImageToPPTConverter()
        converter.process_folder(INPUT_FOLDER, OUTPUT_PPT)