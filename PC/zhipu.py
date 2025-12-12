import os
import cv2
import numpy as np
import base64
import json
import glob
from zhipuai import ZhipuAI
from simple_lama_inpainting import SimpleLama
from paddleocr import PaddleOCR
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from PIL import Image

# ================= 配置区域 =================
# 请在此填入你的 Key
ZHIPU_API_KEY = "38a24dd463fb4002ab9113e2d5f6033c.Pu6ddxsZvONOJR6J" 
INPUT_FOLDER = "images"
OUTPUT_PPT = "glm_final_tuned.pptx"
DEBUG_FOLDER = "debug_output"

# 【关键修改】字号视觉补偿系数
# 之前是 1.3 (偏大) -> 现在改为 0.9 (微缩)
# 如果还觉得大，可以改为 0.8；如果觉得小，改为 1.0
FONT_SIZE_COMPENSATION = 0.9 

# 字体映射表
FONT_MAPPING = {
    "宋体": "SimSun", "SongTi": "SimSun", "SimSun": "SimSun",
    "黑体": "SimHei", "HeiTi": "SimHei", "SimHei": "SimHei",
    "楷体": "KaiTi", "KaiTi": "KaiTi", 
    "仿宋": "FangSong", "FangSong": "FangSong",
    "微软雅黑": "Microsoft YaHei", "YaHei": "Microsoft YaHei",
    "Arial": "Arial", "Times New Roman": "Times New Roman"
}
# ===========================================

class SmartPPTGenerator:
    def __init__(self):
        print("[-] 初始化 GLM 客户端...")
        self.client = ZhipuAI(api_key=ZHIPU_API_KEY)
        
        print("[-] 初始化 OCR 和 LaMa...")
        self.ocr = PaddleOCR(use_angle_cls=True, lang="ch", show_log=False)
        self.lama = SimpleLama()
        
        if not os.path.exists(DEBUG_FOLDER):
            os.makedirs(DEBUG_FOLDER)

    def encode_image_base64(self, cv2_img):
        _, buffer = cv2.imencode('.jpg', cv2_img)
        return base64.b64encode(buffer).decode('utf-8')

    def ask_glm_style(self, crop_img, text_hint):
        base64_str = self.encode_image_base64(crop_img)
        prompt = f"""
        图片包含文字: "{text_hint}"。
        请分析这段文字的视觉样式。注意：在颜色方面，要的是文字颜色，而不是背景颜色
        请严格以JSON格式返回以下三个字段，不要包含Markdown：
        1. "font": 字体最接近哪种？(宋体, 黑体, 楷体, 仿宋, 微软雅黑, Arial)
        2. "color": 文字颜色的十六进制代码 (如 #000000)
        3. "is_bold": 是否为粗体 (true/false)
        """
        try:
            response = self.client.chat.completions.create(
                model="glm-4.6v-flash", 
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": prompt},
                            {"type": "image_url", "image_url": {"url": base64_str}}
                        ]
                    }
                ],
                temperature=0.1, top_p=0.1, max_tokens=1024
            )
            content = response.choices[0].message.content
            content = content.replace("```json", "").replace("```", "").strip()
            return json.loads(content)
        except Exception as e:
            print(f"    [!] GLM 识别出错: {e}")
            return {"font": "微软雅黑", "color": "#000000", "is_bold": False}

    def process_folder(self, folder, output):
        # 修复文件匹配
        extensions = ['*.jpg', '*.jpeg', '*.png', '*.JPG', '*.PNG']
        image_files = []
        for ext in extensions:
            image_files.extend(glob.glob(os.path.join(folder, ext)))
        image_files = sorted(list(set(image_files)))

        if not image_files:
            print(f"[!] 在 {folder} 中未找到图片。")
            return

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for img_path in image_files:
            print(f"\n>>> 正在处理: {os.path.basename(img_path)}")
            self.process_single_image(prs, img_path)

        prs.save(output)
        print(f"\n[√] PPT已生成: {output}")

    def process_single_image(self, prs, img_path):
        img_bgr = cv2.imdecode(np.fromfile(img_path, dtype=np.uint8), cv2.IMREAD_COLOR)
        if img_bgr is None: return
        h_img, w_img = img_bgr.shape[:2]
        
        file_basename = os.path.basename(img_path).split('.')[0]
        
        # 1. OCR
        result = self.ocr.ocr(img_bgr, cls=True)
        valid_blocks = []
        mask = np.zeros((h_img, w_img), dtype=np.uint8)

        if result and result[0]:
            print(f"    [OCR] 找到 {len(result[0])} 个文本块")
            for idx, line in enumerate(result[0]):
                coords = line[0]
                text = line[1][0]
                
                coords_np = np.array(coords, dtype=np.int32)
                x, y, rw, rh = cv2.boundingRect(coords_np)
                
                if rw < 10 or rh < 10: continue

                # 调试保存裁剪图
                pad = 5
                crop = img_bgr[max(0, y-pad):min(h_img, y+rh+pad), max(0, x-pad):min(w_img, x+rw+pad)]
                if crop.size > 0:
                    safe_text = "".join([c for c in text if c.isalnum()])[:10]
                    cv2.imwrite(os.path.join(DEBUG_FOLDER, f"crop_{file_basename}_{idx}_{safe_text}.jpg"), crop)
                    
                    # GLM
                    print(f"      [{idx}] 分析样式: {text[:10]}...", end="\r")
                    style = self.ask_glm_style(crop, text)
                    
                    valid_blocks.append({
                        "text": text,
                        "rect": (x, y, rw, rh),
                        "style": style
                    })
                    cv2.fillPoly(mask, [coords_np], 255)
            print("")

        # 2. LaMa
        if valid_blocks:
            print("    [LaMa] 重绘背景...")
            kernel = np.ones((6, 6), np.uint8)
            mask_dilated = cv2.dilate(mask, kernel, iterations=2)
            img_rgb = cv2.cvtColor(img_bgr, cv2.COLOR_BGR2RGB)
            clean_pil = self.lama(Image.fromarray(img_rgb), Image.fromarray(mask_dilated))
            clean_bgr = cv2.cvtColor(np.array(clean_pil), cv2.COLOR_RGB2BGR)
        else:
            clean_bgr = img_bgr

        # 3. 生成 PPT
        temp_bg = "temp_tune_bg.png"
        cv2.imencode('.png', clean_bgr)[1].tofile(temp_bg)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(temp_bg, 0, 0, width=prs.slide_width, height=prs.slide_height)

        # 4. 回填文字
        scale_x = prs.slide_width / w_img
        scale_y = prs.slide_height / h_img

        for block in valid_blocks:
            x, y, rw, rh = block['rect']
            style = block['style']
            
            left = int(x * scale_x)
            top = int(y * scale_y)
            width = int(rw * scale_x)
            height = int(rh * scale_y)
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            tf = textbox.text_frame
            tf.word_wrap = False
            tf.margin_top = 0; tf.margin_bottom = 0; tf.margin_left = 0
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            
            p = tf.paragraphs[0]
            p.text = block['text']
            
            # 颜色
            hex_color = style.get("color", "#000000").lstrip('#')
            try:
                if len(hex_color) == 3: hex_color = "".join([c*2 for c in hex_color])
                r, g, b = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
                p.font.color.rgb = RGBColor(r, g, b)
            except:
                p.font.color.rgb = RGBColor(0, 0, 0)
            
            # 字体
            p.font.bold = style.get("is_bold", False)
            font_cn = style.get("font", "微软雅黑")
            p.font.name = FONT_MAPPING.get(font_cn, "Microsoft YaHei")
            
            # 【核心修正】字号计算
            # 逻辑：框高(px) -> PPT高度(EMU) -> 磅值(Pt) -> 系数(0.9)
            bbox_height_pt = (rh * scale_y) / 12700
            final_size = bbox_height_pt * FONT_SIZE_COMPENSATION
            
            # 限制最小值 8pt，最大值 80pt (防止极其巨大的字)
            p.font.size = Pt(max(8, min(80, final_size)))

        if os.path.exists(temp_bg):
            os.remove(temp_bg)

if __name__ == "__main__":
    if not os.path.exists(INPUT_FOLDER):
        os.makedirs(INPUT_FOLDER)
        print(f"[!] 请创建文件夹 {INPUT_FOLDER} 并放入图片")
    else:
        generator = SmartPPTGenerator()
        generator.process_folder(INPUT_FOLDER, OUTPUT_PPT)