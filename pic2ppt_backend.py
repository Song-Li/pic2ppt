import os
import cv2
import numpy as np
import glob
import json
from PIL import Image
import torch
from simple_lama_inpainting import SimpleLama
from paddleocr import PaddleOCR
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_VERTICAL_ANCHOR

# ================= 顶级配置 (Server模型配置) =================
MODEL_DIR = "models_server"
DET_DIR = os.path.join(MODEL_DIR, "ch_PP-OCRv4_det_server_infer")
REC_DIR = os.path.join(MODEL_DIR, "ch_PP-OCRv4_rec_server_infer")
CLS_DIR = os.path.join(MODEL_DIR, "ch_ppocr_mobile_v2.0_cls_infer")

# PPT/排版配置
FONT_SIZE_FACTOR = 0.95      # 字体大小系数
MASK_DILATE_ITER = 2         
MASK_PADDING_PIXELS = 2     # Mask 相对于红框的额外安全边距 (用于清理)
PADDING_PIXELS_GREEN = 2     # 绿框的固定外扩像素 
MIN_SCORE = 0.70

class ImageToPPTConverter:
    def __init__(self):
        # 1. 检查模型和初始化 OCR (GPU加速)
        if not os.path.exists(DET_DIR):
            raise FileNotFoundError(f"找不到模型文件: {DET_DIR}\n请先运行 download_models.py")
            
        print("[-] 加载 PaddleOCR Server 模型 (GPU加速中)...")
        self.ocr = PaddleOCR(
            text_detection_model_dir=DET_DIR,
            text_recognition_model_dir=REC_DIR,
            textline_orientation_model_dir=CLS_DIR,
            use_textline_orientation=True,
            device="gpu" 
        )
        
        # 2. 初始化 LaMa
        print("[-] 加载 LaMa 修复模型 (GPU加速中)...")
        self.lama = SimpleLama()
        print("[√] 初始化完成")

    def smart_refine_box(self, img, x, y, w, h):
        """绿框逻辑：返回精准贴合文字的尺寸"""
        if w <= 0 or h <= 0: return x, y, w, h
        
        pad_probe = 2 
        H, W = img.shape[:2]
        y1, y2 = max(0, y - pad_probe), min(H, y + h + pad_probe)
        x1, x2 = max(0, x - pad_probe), min(W, x + w + pad_probe)
        roi = img[y1:y2, x1:x2]
        
        if roi.size == 0: return x, y, w, h

        gray = cv2.cvtColor(roi, cv2.COLOR_BGR2GRAY)
        blurred = cv2.GaussianBlur(gray, (3, 3), 0)
        _, binary = cv2.threshold(blurred, 0, 255, cv2.THRESH_BINARY_INV + cv2.THRESH_OTSU)
        points = cv2.findNonZero(binary)
        
        if points is not None:
            rx, ry, rw, rh = cv2.boundingRect(points)
            
            real_x = x1 + rx
            real_y = y1 + ry
            
            final_x = max(0, real_x - PADDING_PIXELS_GREEN)
            final_y = max(0, real_y - PADDING_PIXELS_GREEN)
            final_w = min(W - final_x, rw + 2*PADDING_PIXELS_GREEN)
            final_h = min(H - final_y, rh + 2*PADDING_PIXELS_GREEN)
            
            return int(final_x), int(final_y), int(final_w), int(final_h)
        return x, y, w, h

    def analyze_style(self, img, box):
        """颜色和粗体分析：基于 HSV 空间，提取最饱和的颜色作为字体色"""
        x, y, w, h = box
        if w <= 0 or h <= 0: return (0,0,0), False
        roi = img[y:y+h, x:x+w] 
        if roi.size == 0: return (0,0,0), False

        is_bold = False
        try:
            roi_hsv = cv2.cvtColor(roi, cv2.COLOR_BGR2HSV)
            lower_bound = np.array([0, 0, 50]) 
            upper_bound = np.array([179, 255, 200]) 
            hsv_mask = cv2.inRange(roi_hsv, lower_bound, upper_bound)
            
            valid_pixels = roi[hsv_mask != 0]
            
            if valid_pixels.size > 0:
                valid_pixels_hsv = cv2.cvtColor(valid_pixels.reshape(1, -1, 3).astype(np.uint8), cv2.COLOR_BGR2HSV).reshape(-1, 3)
                max_saturation_index = np.argmax(valid_pixels_hsv[:, 1])
                color_bgr = valid_pixels[max_saturation_index]
            else:
                color_bgr = cv2.mean(roi)[:3]

            text_color = (int(color_bgr[2]), int(color_bgr[1]), int(color_bgr[0]))
            
            # 粗体判断
            gray = cv2.cvtColor(roi, cv2.COLOR_BGR2GRAY)
            _, binary = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY_INV + cv2.THRESH_OTSU)
            if (cv2.countNonZero(binary) / binary.size) > 0.45:
                is_bold = True
        except: 
            text_color = (0,0,0)
        
        return text_color, is_bold

    def is_noise_or_icon(self, w, h, text, score):
        """强力过滤图标和噪点"""
        if w < 10 or h < 10: return True
        ratio = w / h
        if 0.8 < ratio < 1.3 and max(w, h) < 60 and len(text) < 2:
            if '\u4e00' <= text <= '\u9fff': return False 
            return True
        if len(text) == 1 and not ('\u4e00' <= text <= '\u9fff' or text.isalnum()):
            return True
        return False

    def get_initial_boxes(self, img_path):
        """
        供 GUI 调用的接口：运行 OCR 并返回初始的框列表
        返回格式: [[x, y, w, h, text, raw_x, raw_y, raw_w, raw_h], ...]
        """
        img_bgr = cv2.imdecode(np.fromfile(img_path, dtype=np.uint8), cv2.IMREAD_COLOR)
        if img_bgr is None: return []

        result = self.ocr.ocr(img_bgr)
        lines = result[0] if (result and result[0]) else []
        initial_boxes = []

        for line in lines:
            if not line: continue
            coords = line[0]
            text_content = line[1][0]
            score = line[1][1]
            
            if score < MIN_SCORE: continue

            coords_np = np.array(coords, dtype=np.int32)
            raw_x, raw_y, raw_w, raw_h = cv2.boundingRect(coords_np)
            
            # 1. 过滤图标
            if self.is_noise_or_icon(raw_w, raw_h, text_content, score):
                continue
            
            # 2. 计算绿框 (精修尺寸)
            g_x, g_y, g_rw, g_rh = self.smart_refine_box(img_bgr, raw_x, raw_y, raw_w, raw_h)
            
            if g_rw < 5 or g_rh < 5: continue
            
            # 返回绿框作为初始可编辑的框
            initial_boxes.append([g_x, g_y, g_rw, g_rh, text_content, raw_x, raw_y, raw_w, raw_h])

        return initial_boxes

    def generate_ppt_from_manual_data(self, image_files, output_path, manual_data_dir="manual_output"):
        """
        根据手动调整后的JSON数据生成最终PPT
        """
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        for img_path in image_files:
            filename = os.path.basename(img_path)
            json_filename = os.path.splitext(filename)[0] + "_manual.json"
            json_path = os.path.join(manual_data_dir, json_filename)
            
            # 检查是否有手动数据
            if not os.path.exists(json_path):
                print(f"跳过 {filename}: 未找到手动调整数据 ({json_filename})")
                continue
            
            print(f"[-] 正在处理 {filename} (使用手动框)...")
            
            with open(json_path, 'r', encoding='utf-8') as f:
                manual_data = json.load(f)
            
            img_bgr = cv2.imdecode(np.fromfile(img_path, dtype=np.uint8), cv2.IMREAD_COLOR)
            if img_bgr is None: continue
            h_img, w_img = img_bgr.shape[:2]

            valid_blocks = []
            mask = np.zeros((h_img, w_img), dtype=np.uint8)

            for block in manual_data["boxes"]:
                x, y, rw, rh = block["x"], block["y"], block["w"], block["h"]
                text_content = block["text"]
                
                # 1. 颜色和样式分析 (基于手动调整后的绿框尺寸)
                color, is_bold = self.analyze_style(img_bgr, (x, y, rw, rh)) 
                
                valid_blocks.append({
                    "text": text_content,
                    "rect": (x, y, rw, rh),
                    "color": color,
                    "bold": is_bold
                })
                
                # 2. Mask 绘制：【修复点】完全依赖手动调整后的框 + 安全边距
                m_x1 = max(0, x - MASK_PADDING_PIXELS)
                m_y1 = max(0, y - MASK_PADDING_PIXELS)
                m_x2 = min(w_img, x + rw + MASK_PADDING_PIXELS)
                m_y2 = min(h_img, y + rh + MASK_PADDING_PIXELS)
                cv2.rectangle(mask, (m_x1, m_y1), (m_x2, m_y2), 255, -1)


            # 3. LaMa 修复
            if valid_blocks:
                kernel = np.ones((3, 3), np.uint8)
                mask_dilated = cv2.dilate(mask, kernel, iterations=MASK_DILATE_ITER)
                
                img_rgb = cv2.cvtColor(img_bgr, cv2.COLOR_BGR2RGB)
                try:
                    clean_pil = self.lama(Image.fromarray(img_rgb), Image.fromarray(mask_dilated))
                    clean_bgr = cv2.cvtColor(np.array(clean_pil), cv2.COLOR_RGB2BGR)
                except Exception as e:
                    print(f"   [!] LaMa Error: {e}")
                    clean_bgr = img_bgr
            else:
                clean_bgr = img_bgr

            # 4. 插入 PPT
            temp_bg = f"temp_{os.path.splitext(filename)[0]}.png"
            cv2.imencode('.png', clean_bgr)[1].tofile(temp_bg)
            
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(temp_bg, 0, 0, width=prs.slide_width, height=prs.slide_height)
            
            scale_x = prs.slide_width / w_img
            scale_y = prs.slide_height / h_img
            
            for block in valid_blocks:
                x, y, rw, rh = block['rect']
                text_content = block['text']
                
                textbox = slide.shapes.add_textbox(
                    int(x * scale_x), int(y * scale_y), 
                    int(rw * scale_x), int(rh * scale_y)
                )
                tf = textbox.text_frame
                tf.word_wrap = False 
                tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0
                
                p = tf.paragraphs[0]
                p.text = text_content
                p.font.color.rgb = RGBColor(*block['color'])
                p.font.bold = block['bold']
                p.font.name = "微软雅黑"
                
                # --- 字号双重约束 ---
                H_pt = (rh / h_img) * (prs.slide_height.inches * 72)
                char_count = len(text_content)
                W_pt = (rw / w_img) * (prs.slide_width.inches * 72)
                W_pt_constraint = W_pt / char_count / 1.0 
                
                final_pt_size = min(H_pt, W_pt_constraint)
                
                p.font.size = Pt(max(9, final_pt_size * FONT_SIZE_FACTOR))

                if os.path.exists(temp_bg): os.remove(temp_bg)
        
        # 核心：只有在循环结束后，prs.save() 才能保证文件创建成功
        try:
            prs.save(output_path)
            print("\n[√] PPT生成完成！")
        except Exception as e:
            print(f"\n[!] PPT保存失败: {e}")
        
        print("\n[√] PPT生成完成！")