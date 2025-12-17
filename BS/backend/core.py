import os
import cv2
import numpy as np
import glob
import json
from PIL import Image
import torch
from simple_lama_inpainting import SimpleLama
from paddleocr import PaddleOCR
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_VERTICAL_ANCHOR

# ================= 顶级配置 (v4 Server 模型) =================
# 请确保这里指向你下载的模型文件夹
MODEL_BASE = "./models"
DET_DIR = os.path.join(MODEL_BASE, "ch_PP-OCRv4_det_server_infer")
REC_DIR = os.path.join(MODEL_BASE, "ch_PP-OCRv4_rec_server_infer")
CLS_DIR = os.path.join(MODEL_BASE, "ch_ppocr_mobile_v2.0_cls_infer")

# PPT/排版配置
FONT_SIZE_FACTOR = 0.95      
MASK_DILATE_ITER = 2         
MASK_PADDING_PIXELS = 2     
PADDING_PIXELS_GREEN = 2     
MIN_SCORE = 0.70

class ImageProcessor:
    def __init__(self):
        # 1. 检查模型文件是否存在
        if not os.path.exists(DET_DIR) or not os.path.exists(REC_DIR):
            print(f"[!] 错误：找不到 Server 模型文件！")
            print(f"    请先运行 'download_models.py' 下载模型到 {MODEL_BASE} 文件夹。")
            raise FileNotFoundError("本地 Server 模型未找到")

        print("[-] 加载 PaddleOCR v4 Server 模型 (GPU加速中)...")
        # 显式指定本地模型路径
        self.ocr = PaddleOCR(
            text_detection_model_dir=DET_DIR,       
            text_recognition_model_dir=REC_DIR,     
            textline_orientation_model_dir=CLS_DIR, 
            use_angle_cls=True,                     
            lang="ch",
        )
        
        # 2. 初始化 LaMa
        print("[-] 加载 LaMa 修复模型 (GPU加速中)...")
        self.lama = SimpleLama()
        if torch.cuda.is_available():
            self.lama.device = torch.device("cuda")
            if hasattr(self.lama, 'model'):
                self.lama.model.to(self.lama.device)
        
        print("[√] 初始化完成 (v4 Server版)")

    def smart_refine_box(self, img, x, y, w, h):
        """绿框逻辑：返回精准贴合文字的尺寸"""
        if w <= 0 or h <= 0: return x, y, w, h
        
        pad_probe = 2 
        H, W = img.shape[:2]
        y1, y2 = max(0, y - pad_probe), min(H, y + h + pad_probe)
        x1, x2 = max(0, x - pad_probe), min(W, x + w + pad_probe)
        roi = img[y1:y2, x1:x2]
        
        if roi.size == 0: return x, y, w, h

        gray = cv2.cvtColor(roi, cv2.COLOR_BGR2GRAY)
        blurred = cv2.GaussianBlur(gray, (3, 3), 0)
        _, binary = cv2.threshold(blurred, 0, 255, cv2.THRESH_BINARY_INV + cv2.THRESH_OTSU)
        points = cv2.findNonZero(binary)
        
        if points is not None:
            rx, ry, rw, rh = cv2.boundingRect(points)
            final_x = max(0, x1 + rx - PADDING_PIXELS_GREEN)
            final_y = max(0, y1 + ry - PADDING_PIXELS_GREEN)
            final_w = min(W - final_x, rw + 2*PADDING_PIXELS_GREEN)
            final_h = min(H - final_y, rh + 2*PADDING_PIXELS_GREEN)
            return int(final_x), int(final_y), int(final_w), int(final_h)
        return x, y, w, h

    def analyze_style(self, img, box):
        """颜色和粗体分析"""
        x, y, w, h = box
        if w <= 0 or h <= 0: return (0,0,0), False
        roi = img[y:y+h, x:x+w] 
        if roi.size == 0: return (0,0,0), False

        is_bold = False
        try:
            roi_hsv = cv2.cvtColor(roi, cv2.COLOR_BGR2HSV)
            mask = cv2.inRange(roi_hsv, np.array([0, 0, 50]), np.array([179, 255, 200]))
            valid_pixels = roi[mask != 0]
            
            if valid_pixels.size > 0:
                valid_hsv = cv2.cvtColor(valid_pixels.reshape(1, -1, 3).astype(np.uint8), cv2.COLOR_BGR2HSV).reshape(-1, 3)
                idx = np.argmax(valid_hsv[:, 1])
                color_bgr = valid_pixels[idx]
            else:
                color_bgr = cv2.mean(roi)[:3]

            text_color = (int(color_bgr[2]), int(color_bgr[1]), int(color_bgr[0]))
            
            gray = cv2.cvtColor(roi, cv2.COLOR_BGR2GRAY)
            _, binary = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY_INV + cv2.THRESH_OTSU)
            if (cv2.countNonZero(binary) / binary.size) > 0.45:
                is_bold = True
        except: 
            text_color = (0,0,0)
        
        return text_color, is_bold

    def is_noise_or_icon(self, w, h, text, score):
        if w < 10 or h < 10: return True
        ratio = w / h
        if 0.8 < ratio < 1.3 and max(w, h) < 60 and len(text) < 2:
            if '\u4e00' <= text <= '\u9fff': return False 
            return True
        if len(text) == 1 and not ('\u4e00' <= text <= '\u9fff' or text.isalnum()):
            return True
        return False

    def get_initial_boxes(self, img_path):
        """执行 OCR 并返回原始框数据"""
        img_bgr = cv2.imdecode(np.fromfile(img_path, dtype=np.uint8), cv2.IMREAD_COLOR)
        if img_bgr is None: return [], None

        result = self.ocr.ocr(img_bgr)
        lines = result[0] if (result and result[0]) else []
        initial_boxes = []

        for line in lines:
            if not line: continue
            coords = line[0]
            text_content = line[1][0]
            score = line[1][1]
            
            if score < MIN_SCORE: continue

            coords_np = np.array(coords, dtype=np.int32)
            raw_x, raw_y, raw_w, raw_h = cv2.boundingRect(coords_np)
            
            if self.is_noise_or_icon(raw_w, raw_h, text_content, score):
                continue
            
            g_x, g_y, g_rw, g_rh = self.smart_refine_box(img_bgr, raw_x, raw_y, raw_w, raw_h)
            
            if g_rw < 5 or g_rh < 5: continue
            # 返回: [精修x, 精修y, 精修w, 精修h, 文本]
            initial_boxes.append([g_x, g_y, g_rw, g_rh, text_content])

        return initial_boxes, img_bgr

    # ================= 新增：对接 main.py 的接口 =================

    def process_image(self, img_path):
        """处理上传图片：OCR -> 分析样式 -> 返回前端 JSON"""
        print(f"[-] 正在处理图片: {img_path}")
        
        # 1. 获取 OCR 框和原图
        raw_boxes, img_bgr = self.get_initial_boxes(img_path)
        
        if img_bgr is None:
            raise ValueError("图片加载失败")
            
        h, w = img_bgr.shape[:2]

        formatted_boxes = []
        
        # 2. 遍历框，补充颜色/粗体信息，格式化为字典
        for idx, box in enumerate(raw_boxes):
            bx, by, bw, bh, text = box
            color, is_bold = self.analyze_style(img_bgr, (bx, by, bw, bh))
            
            formatted_boxes.append({
                "id": idx,
                "x": int(bx), "y": int(by),
                "w": int(bw), "h": int(bh),
                "text": text,
                "color": color,   # (R, G, B)
                "is_bold": is_bold
            })

        # 3. 返回给 main.py
        return {
            "width": w,
            "height": h,
            "boxes": formatted_boxes
        }

    def generate_ppt(self, img_path, boxes, output_path):
        """接收前端确认后的框，执行修图并生成 PPT"""
        print(f"[-] 开始生成 PPT: {output_path}")
        
        img_bgr = cv2.imdecode(np.fromfile(img_path, dtype=np.uint8), cv2.IMREAD_COLOR)
        if img_bgr is None: raise ValueError("无法读取图片")
        
        h_img, w_img = img_bgr.shape[:2]
        
        # 1. 准备 LaMa 遮罩 (Mask)
        mask = np.zeros((h_img, w_img), dtype=np.uint8)
        
        for box in boxes:
            # 前端传回来的坐标
            x, y, w, h = box['x'], box['y'], box['w'], box['h']
            
            # 为了修图干净，稍微扩大遮罩
            m_x1 = max(0, x - MASK_PADDING_PIXELS)
            m_y1 = max(0, y - MASK_PADDING_PIXELS)
            m_x2 = min(w_img, x + w + MASK_PADDING_PIXELS)
            m_y2 = min(h_img, y + h + MASK_PADDING_PIXELS)
            
            cv2.rectangle(mask, (m_x1, m_y1), (m_x2, m_y2), 255, -1)

        # 2. 执行 LaMa 修复 (Inpainting)
        if len(boxes) > 0:
            kernel = np.ones((3, 3), np.uint8)
            mask_dilated = cv2.dilate(mask, kernel, iterations=MASK_DILATE_ITER)
            img_rgb = cv2.cvtColor(img_bgr, cv2.COLOR_BGR2RGB)
            
            try:
                # 调用 LaMa
                clean_pil = self.lama(Image.fromarray(img_rgb), Image.fromarray(mask_dilated))
                clean_bgr = cv2.cvtColor(np.array(clean_pil), cv2.COLOR_RGB2BGR)
            except Exception as e:
                print(f"   [!] LaMa Error: {e}")
                clean_bgr = img_bgr
        else:
            clean_bgr = img_bgr

        # 3. 创建 PPT
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # 空白版式

        # 保存临时背景图并插入
        temp_bg = output_path.replace(".pptx", "_temp.png")
        cv2.imencode('.png', clean_bgr)[1].tofile(temp_bg)
        slide.shapes.add_picture(temp_bg, 0, 0, width=prs.slide_width, height=prs.slide_height)

        # 4. 插入可编辑文本框
        scale_x = prs.slide_width / w_img
        scale_y = prs.slide_height / h_img
        
        for box in boxes:
            x, y, w, h = box['x'], box['y'], box['w'], box['h']
            text = box['text']
            # 如果前端传回了颜色，使用前端的；否则默认黑色
            color = box.get('color', [0,0,0]) 
            
            textbox = slide.shapes.add_textbox(
                int(x * scale_x), int(y * scale_y), 
                int(w * scale_x), int(h * scale_y)
            )
            tf = textbox.text_frame
            tf.word_wrap = False 
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            # 消除默认边距
            tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0
            
            p = tf.paragraphs[0]
            p.text = text
            # 设置颜色
            p.font.color.rgb = RGBColor(int(color[0]), int(color[1]), int(color[2]))
            p.font.name = "微软雅黑"
            
            # 动态计算字号
            H_pt = (h / h_img) * (prs.slide_height.inches * 72)
            char_count = len(text) if len(text) > 0 else 1
            W_pt = (w / w_img) * (prs.slide_width.inches * 72)
            W_pt_constraint = W_pt / char_count / 1.0 
            final_pt_size = min(H_pt, W_pt_constraint)
            p.font.size = Pt(max(9, final_pt_size * FONT_SIZE_FACTOR))

        prs.save(output_path)
        
        # 清理临时背景文件
        if os.path.exists(temp_bg):
            os.remove(temp_bg)
        
        print("[√] PPT 生成完毕")

if __name__ == "__main__":
    try:
        processor = ImageProcessor()
        print("核心模块加载成功。请运行 main.py 启动服务器。")
    except Exception as e:
        print(f"初始化失败: {e}")