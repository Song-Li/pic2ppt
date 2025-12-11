import os
import cv2
import numpy as np
import glob
from PIL import Image
from simple_lama_inpainting import SimpleLama
from paddleocr import PaddleOCR
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR

# ================= 配置区域 =================
INPUT_FOLDER = "images"          # 输入文件夹（请确保此文件夹存在并放入图片）
OUTPUT_PPT = "result_perfect.pptx" # 输出文件名
MIN_SCORE = 0.60                 # OCR置信度过滤
FILTER_ICON_RATIO = 1.2          # 图标过滤阈值

# 字号微调系数
# 0.95 是一个比较平衡的值，既能填满框，又不会因为Ascender/Descender溢出
FONT_SIZE_FACTOR = 0.95          
# ===========================================

class ImageToPPTConverter:
    def __init__(self):
        print("[-] 初始化 PaddleOCR (启用角度分类)...")
        # 如果下载模型慢，可以手动下载模型放本地
        self.ocr = PaddleOCR(use_angle_cls=True, lang="ch", show_log=False)
        print("[-] 初始化 LaMa 修复模型...")
        self.lama = SimpleLama()
        print("[√] 初始化完成")

    def analyze_style(self, img, box):
        """
        样式分析：提取颜色和判断是否加粗
        """
        x, y, w, h = box
        if w <= 0 or h <= 0: return (0,0,0), False
        
        # 裁剪出文字区域，稍微扩大一点以包含背景对比
        pad = 2
        y1, y2 = max(0, y-pad), min(img.shape[0], y+h+pad)
        x1, x2 = max(0, x-pad), min(img.shape[1], x+w+pad)
        roi = img[y1:y2, x1:x2]
        
        if roi.size == 0: return (0, 0, 0), False

        # --- 1. 颜色提取 (K-Means) ---
        # 尝试聚类背景和前景，提取较少的那部分作为文字色
        data = roi.reshape((-1, 3))
        data = np.float32(data)
        criteria = (cv2.TERM_CRITERIA_EPS + cv2.TERM_CRITERIA_MAX_ITER, 10, 1.0)
        
        try:
            K = 2
            ret, label, center = cv2.kmeans(data, K, None, criteria, 10, cv2.KMEANS_RANDOM_CENTERS)
            counts = np.bincount(label.flatten())
            # 像素少的是文字
            if len(counts) >= 2:
                text_cluster_index = np.argmin(counts)
                color_bgr = center[text_cluster_index]
            else:
                color_bgr = center[0]
        except:
            color_bgr = cv2.mean(roi)[:3] # 失败回落到均值

        text_color = (int(color_bgr[2]), int(color_bgr[1]), int(color_bgr[0])) # BGR -> RGB

        # --- 2. 粗体检测 ---
        is_bold = False
        try:
            gray = cv2.cvtColor(roi, cv2.COLOR_BGR2GRAY)
            # Otsu 二值化
            _, binary = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY_INV + cv2.THRESH_OTSU)
            pixel_ratio = cv2.countNonZero(binary) / (binary.shape[0] * binary.shape[1])
            # 笔画占比高且高度足够，认为是粗体
            if pixel_ratio > 0.35 and h > 15: 
                is_bold = True
        except:
            pass

        return text_color, is_bold

    def is_likely_icon(self, w, h, text):
        """过滤掉非文字的小图标"""
        ratio = max(w, h) / min(w, h)
        if w < 50 and h < 50 and ratio < FILTER_ICON_RATIO and len(text) < 2:
            return True
        return False

    def process_folder(self, folder_path, output_path):
        if not os.path.exists(folder_path):
            print(f"[!] 错误：文件夹 '{folder_path}' 不存在")
            return

        image_files = sorted(glob.glob(os.path.join(folder_path, "*.[jpg][png][jpeg]*")))
        
        if not image_files:
            print("[!] 未找到图片，请检查 input 文件夹")
            return

        prs = Presentation()
        # 设置 PPT 为 16:9 宽屏
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for i, img_path in enumerate(image_files):
            print(f">>> [{i+1}/{len(image_files)}] 处理: {os.path.basename(img_path)}")
            self.process_one_slide(prs, img_path)

        prs.save(output_path)
        print(f"\n[√] 全部完成！PPT已保存为: {output_path}")

    def process_one_slide(self, prs, img_path):
        # 读取图片
        img_bgr = cv2.imdecode(np.fromfile(img_path, dtype=np.uint8), cv2.IMREAD_COLOR)
        if img_bgr is None:
            print(f"   [!] 无法读取图片: {img_path}")
            return
            
        h_img, w_img = img_bgr.shape[:2]

        # 1. OCR 识别
        result = self.ocr.ocr(img_bgr, cls=True)
        valid_blocks = []
        mask = np.zeros((h_img, w_img), dtype=np.uint8)

        if result and result[0]:
            for line in result[0]:
                coords = line[0]
                text = line[1][0]
                score = line[1][1]

                if score < MIN_SCORE: continue
                
                coords_np = np.array(coords, dtype=np.int32)
                x, y, rw, rh = cv2.boundingRect(coords_np)
                
                if self.is_likely_icon(rw, rh, text): continue

                # 分析样式
                color, is_bold = self.analyze_style(img_bgr, (x, y, rw, rh))
                
                valid_blocks.append({
                    "text": text,
                    "rect": (x, y, rw, rh),
                    "color": color,
                    "bold": is_bold
                })
                
                # 在 mask 上把文字区域涂白，用于后续擦除
                cv2.fillPoly(mask, [coords_np], 255)

        # 2. 图片修复 (LaMa Inpainting)
        if valid_blocks:
            # 膨胀 mask 以覆盖文字边缘
            kernel = np.ones((6, 6), np.uint8)
            mask_dilated = cv2.dilate(mask, kernel, iterations=2)
            
            img_rgb = cv2.cvtColor(img_bgr, cv2.COLOR_BGR2RGB)
            try:
                # 传入 PIL Image
                clean_pil = self.lama(Image.fromarray(img_rgb), Image.fromarray(mask_dilated))
                clean_bgr = cv2.cvtColor(np.array(clean_pil), cv2.COLOR_RGB2BGR)
            except Exception as e:
                print(f"   [!] LaMa 修复失败，使用原图背景: {e}")
                clean_bgr = img_bgr
        else:
            clean_bgr = img_bgr

        # 3. 将修复后的图片作为背景插入 PPT
        temp_bg = "temp_bg_processing.png"
        cv2.imencode('.png', clean_bgr)[1].tofile(temp_bg)
        
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # 6 = 空白版式
        slide.shapes.add_picture(temp_bg, 0, 0, width=prs.slide_width, height=prs.slide_height)
        
        # 4. 回填文字
        scale_x = prs.slide_width / w_img
        scale_y = prs.slide_height / h_img
        
        for block in valid_blocks:
            x, y, rw, rh = block['rect']
            
            # 坐标转换 (Pixel -> EMU)
            left = int(x * scale_x)
            top = int(y * scale_y)
            width = int(rw * scale_x)
            height = int(rh * scale_y)
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            tf = textbox.text_frame
            tf.word_wrap = False 
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            # 移除边距，使位置更精准
            tf.margin_top = 0
            tf.margin_bottom = 0
            tf.margin_left = 0
            tf.margin_right = 0
            
            p = tf.paragraphs[0]
            p.text = block['text']
            p.font.color.rgb = RGBColor(*block['color'])
            p.font.bold = block['bold']
            p.font.name = "微软雅黑" # 默认中文字体
            
            # --- 字号核心修复逻辑 ---
            # 1. 计算文字高度在图片中的占比
            height_ratio = rh / h_img
            
            # 2. PPT总高度 (Points) = 英寸 * 72
            slide_total_pt = prs.slide_height.inches * 72
            
            # 3. 基础字号 = 占比 * PPT总高度
            base_pt = height_ratio * slide_total_pt
            
            # 4. 应用微调系数
            final_pt = base_pt * FONT_SIZE_FACTOR
            
            # 5. 设置字号 (限制最小为10号字，防止过小)
            p.font.size = Pt(max(10, final_pt))

        # 清理临时文件
        if os.path.exists(temp_bg):
            os.remove(temp_bg)

if __name__ == "__main__":
    if not os.path.exists(INPUT_FOLDER):
        os.makedirs(INPUT_FOLDER)
        print(f"已创建文件夹 '{INPUT_FOLDER}'，请将图片放入其中后再次运行。")
    else:
        converter = ImageToPPTConverter()
        converter.process_folder(INPUT_FOLDER, OUTPUT_PPT)